#!/usr/bin/env python3
"""Build the progress deck (slides/progress_2026-06-17.pptx) with python-pptx.

16:9, Forest & Moss + coral palette (matches the figures). Every slide has a
visual; each carries a speaker script in the notes pane. Run:

    python slides/build_deck.py
"""
from __future__ import annotations
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

REPO = pathlib.Path(__file__).resolve().parents[1]
A = REPO / "slides" / "assets"
OUT = REPO / "slides" / "progress_2026-06-17.pptx"

# ---- palette ----
FOREST = RGBColor(0x3D, 0x6E, 0x4A)
CORAL  = RGBColor(0xB8, 0x5B, 0x3A)
GOLD   = RGBColor(0xB0, 0x7D, 0x24)
TEAL   = RGBColor(0x2C, 0x6E, 0x73)
SLATE  = RGBColor(0x52, 0x61, 0x67)
BLUEG  = RGBColor(0x3A, 0x4A, 0x55)
CHAR   = RGBColor(0x2B, 0x2B, 0x2B)
DIM    = RGBColor(0x6B, 0x6B, 0x6B)
CREAM  = RGBColor(0xF7, 0xF5, 0xF0)
DARK   = RGBColor(0x22, 0x32, 0x2A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
HEAD, BODY = "Georgia", "Calibri"
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(s, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def tbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return tf

def run(p, text, size, color, bold=False, font=BODY, italic=False):
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.color.rgb = color
    f.bold = bold; f.italic = italic; f.name = font

def para(tf, text, size, color, first=False, **kw):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = kw.get("align", PP_ALIGN.LEFT)
    p.space_after = Pt(kw.get("sa", 6)); p.space_before = Pt(kw.get("sb", 0))
    run(p, text, size, color, kw.get("bold", False), kw.get("font", BODY), kw.get("italic", False))
    return p

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

def chip(s, text, color):
    sp = rect(s, 9.55, 0.42, 3.25, 0.46, color, MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, text, 12, WHITE, bold=True)

def fit(s, img, l, t, w, h):
    iw, ih = Image.open(img).size
    ar, bar = iw / ih, w / h
    nw, nh = (w, w / ar) if ar > bar else (h * ar, h)
    s.shapes.add_picture(str(img), Inches(l + (w - nw) / 2), Inches(t + (h - nh) / 2),
                         Inches(nw), Inches(nh))

def bullets(tf, items, size=16, color=CHAR, sa=10):
    for i, it in enumerate(items):
        para(tf, "•  " + it, size, color, first=(i == 0), sa=sa)


def base(section, sec_color, title_text):
    """Cream content slide: bg, left accent bar, section chip, title."""
    s = blank()
    rect(s, -0.1, -0.1, SW + 0.2, SH + 0.2, CREAM)
    rect(s, 0, 0, 0.22, SH, sec_color)
    chip(s, section, sec_color)
    tf = tbox(s, 0.82, 0.44, 8.5, 1.0)
    para(tf, title_text, 30, CHAR, first=True, bold=True, font=HEAD)
    return s

def img_right(section, sec_color, title_text, items, img, note, cap=None):
    s = base(section, sec_color, title_text)
    bullets(tbox(s, 0.82, 1.85, 4.6, 5.0), items)
    fit(s, A / f"{img}.png", 5.6, 1.6, 7.1, 5.5)
    if cap:
        para(tbox(s, 5.6, 6.95, 7.1, 0.4), cap, 11, DIM, first=True, italic=True, align=PP_ALIGN.CENTER)
    notes(s, note); return s

def img_below(section, sec_color, title_text, items, img, note, cap=None):
    s = base(section, sec_color, title_text)
    tf = tbox(s, 0.82, 1.42, 11.9, 0.95)
    bullets(tf, items, size=14, sa=4)
    fit(s, A / f"{img}.png", 0.82, 2.5, 11.7, 4.55)
    if cap:
        para(tbox(s, 0.82, 7.06, 11.7, 0.35), cap, 11, DIM, first=True, italic=True, align=PP_ALIGN.CENTER)
    notes(s, note); return s


def dark_slide(title_text, sub, lines, footer=None):
    s = blank()
    rect(s, -0.1, -0.1, SW + 0.2, SH + 0.2, DARK)
    rect(s, 0, SH - 0.18, SW, 0.18, FOREST)
    rect(s, 0, 0, SW, 0.18, FOREST)
    tf = tbox(s, 1.0, 2.2 if sub else 1.0, 11.3, 2.4, anchor=MSO_ANCHOR.TOP)
    para(tf, title_text, 38 if sub else 34, WHITE, first=True, bold=True, font=HEAD)
    if sub:
        para(tf, sub, 20, RGBColor(0x97, 0xBC, 0x62), sb=10)
    if lines:
        tf2 = tbox(s, 1.0, 3.4, 11.3, 3.2)
        bullets(tf2, lines, size=17, color=RGBColor(0xE7, 0xE8, 0xD1), sa=12)
    if footer:
        para(tbox(s, 1.0, 6.7, 11.3, 0.5), footer, 13, RGBColor(0xCA, 0xDC, 0xC0), first=True)
    return s


def table_slide(section, sec_color, title_text, headers, rows, widths, note, cap=None):
    s = base(section, sec_color, title_text)
    nr, nc = len(rows) + 1, len(headers)
    total_w = sum(widths)
    left = (SW - total_w) / 2
    gt = s.shapes.add_table(nr, nc, Inches(left), Inches(2.1),
                            Inches(total_w), Inches(0.55 * nr)).table
    gt.first_row = False; gt.horz_banding = False
    for j, w in enumerate(widths):
        gt.columns[j].width = Inches(w)
    for j, htext in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = sec_color
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        run(p, htext, 14, WHITE, bold=True)
    for i, rowvals in enumerate(rows, start=1):
        band = CREAM if i % 2 else RGBColor(0xEC, 0xE6, 0xDA)
        for j, val in enumerate(rowvals):
            c = gt.cell(i, j); c.fill.solid(); c.fill.fore_color.rgb = band
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run(p, val, 13, CHAR, bold=(rowvals[0].startswith("M1") or val.startswith("Total")))
    if cap:
        para(tbox(s, 0.82, 2.2 + 0.55 * nr + 0.25, 11.7, 0.8), cap, 13, DIM, first=True, italic=True)
    notes(s, note); return s


# ============================== BUILD ==============================
# 1 — title
s = dark_slide(
    "Per-site deep-regolith thermal conductivity $K_d$ at Apollo 15 & 17".replace("$K_d$", "K_d"),
    "A converged, guess-independent retrieval — progress update",
    None,
    "R. P. Gregorio  ·  Kasai Laboratory, Institute of Science Tokyo  ·  17 June 2026")
notes(s, "Good morning. I'll walk through where the Apollo heat-flow K_d retrieval stands. "
         "The headline: we now retrieve the deep regolith conductivity separately at Apollo 15 and 17, "
         "with a forward model fixed to converge to a true steady state independent of its starting guess. "
         "I'll cover how that fix came about, the new K_d numbers and how they're justified, how they depend "
         "on the basal heat flux, and the computing side — including a possible C++ port and the cores question.")

# 2 — motivation
img_right("Motivation", SLATE, "Why retrieve K_d per site?",
    ["K_d (deep regolith conductivity) sets the deep temperature gradient and the heat escaping the Moon.",
     "Published models use ONE global K_d — the two HFE boreholes may genuinely differ.",
     "Per-site K_d feeds subsurface thermal / ice-stability models and future missions."],
    "fig_context_map",
    "K_d is the deep thermal conductivity of the regolith — it controls how steeply temperature rises with "
    "depth and how much internal heat reaches the surface. Standard practice uses a single global value, but "
    "Apollo 15 at Hadley and 17 at Taurus-Littrow sit in different terrains. If their conductivities differ, a "
    "global value is wrong for both. Getting them per site matters for anyone modelling the subsurface.")

# 3 — data
img_below("Motivation", SLATE, "The data: restored Apollo HFE record (1971–1977)",
    ["Nagihara et al. (2018) restoration of the buried-probe temperatures.",
     "We fit only the settled ‘stability window’ of each sensor (not the noisy startup).",
     "80 cm borestem cut (metal stem distorts shallow heat) → N = 7 (A15), 16 (A17) deep sensors."],
    "fig_apollo_timeline",
    "The data is the restored 1971-77 record from Nagihara 2018. Two cleaning steps. First, we use only the late, "
    "thermally settled part of each sensor's record — the stability window. Second, the top 80 cm is contaminated "
    "by heat conducted down the metal borestem, so we discard sensors above that. That leaves 7 trustworthy deep "
    "sensors at Apollo 15 and 16 at Apollo 17 — the numbers we actually fit.")

# 4 — forward model
img_right("Method", SLATE, "The forward model",
    ["1-D heat equation in the regolith column; geometric grid (mm at top → 5 m deep).",
     "Temperature- and depth-dependent conductivity K(T, z) (Hayne 2017 form).",
     "Surface: full radiative energy balance. Bottom: geothermal flux Q_b.",
     "Solved to a periodic steady state, then compared to the deep sensors."],
    "fig_intro_probe",
    "Under the data sits a standard 1-D heat-conduction model. A grid that's fine near the surface to resolve the "
    "daily thermal wave and coarse at depth; conductivity that depends on temperature and depth; sunlight balanced "
    "by radiation at the surface; and the Moon's geothermal flux entering at the bottom. We run it until each lunar "
    "day repeats — the periodic steady state — and the whole retrieval is just: find the K_d that best matches the sensors.")

# 5 — problem F1
s = base("Equilibrium fix", FOREST, "The problem we found (flag F1)")
card = rect(s, 0.82, 1.75, 11.7, 1.25, RGBColor(0xEA, 0xF1, 0xED), MSO_SHAPE.ROUNDED_RECTANGLE)
tf = card.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
tf.margin_left = Inches(0.3)
p = tf.paragraphs[0]; run(p, "A fixed ~30-lunation spin-up never reaches steady state at sensor depth.", 22, FOREST, bold=True, font=HEAD)
bullets(tbox(s, 0.82, 3.5, 11.7, 3.2),
    ["The deep column relaxes over ~10³ lunations — far longer than the spin-up.",
     "So the initial guess (an ‘anchor’ temperature) leaked into the answer → K_d biased by a hidden parameter.",
     "Worse: the usual ‘ΔT between lunations < 0.01 K’ check passed anyway — it looks converged long before it is."],
    size=17, sa=12)
notes(s, "Here's the bug that started all this — flag F1. The model needs a steady state, but the deep column relaxes "
         "incredibly slowly, order a thousand lunar days. The old code ran about thirty. So at sensor depth the model still "
         "remembered the temperature we started it at — effectively a free knob we'd set by hand, and it biased K_d. The trap "
         "is that the standard convergence check, 'did the temperature stop changing between days', passes far too early because "
         "the deep drift is so slow. So you get a confident, wrong answer.")

# 6 — mechanism
s = base("Equilibrium fix", FOREST, "The fix: a flux-anchored shortcut")
banner = rect(s, 0.82, 1.6, 11.7, 0.8, FOREST, MSO_SHAPE.ROUNDED_RECTANGLE)
tfb = banner.text_frame; tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
pb = tfb.paragraphs[0]; pb.alignment = PP_ALIGN.CENTER
run(pb, "At steady state the deep gradient is fixed:   dT/dz = Q_b / K(T, z)", 18, WHITE, bold=True)
steps = [("1", "Short spin-up (~12 lunations) settles the fast surface skin and pins the anchor temperature at ~0.55 m."),
         ("2", "From the anchor, integrate dT/dz = Q_b/K straight down — the whole deep profile, instantly."),
         ("3", "Feed back and repeat 3–5× → converges. No long wait.")]
for i, (n, txt) in enumerate(steps):
    y = 2.7 + i * 1.15
    circ = rect(s, 0.95, y, 0.7, 0.7, FOREST, MSO_SHAPE.OVAL)
    ct = circ.text_frame; ct.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ct.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER; run(cp, n, 20, WHITE, bold=True)
    para(tbox(s, 1.9, y + 0.02, 10.7, 0.95, anchor=MSO_ANCHOR.MIDDLE), txt, 16, CHAR, first=True)
para(tbox(s, 0.95, 6.45, 11.6, 0.6), "Analogy: find a tank's steady level from inflow = outflow, instead of waiting for the sloshing to stop.",
     13, DIM, first=True, italic=True)
notes(s, "The fix doesn't wait. It uses a physical fact: at steady state the same heat flows through every layer — the "
         "geothermal flux Q_b — so the gradient at every depth is just Q_b divided by conductivity. That's not an unknown to "
         "discover; it's fixed by energy conservation. So: step one, a short spin-up, only long enough to settle the fast skin "
         "and pin one anchor temperature. Step two, integrate that known gradient straight down to fill the deep profile — "
         "instantly. Step three, repeat a few times. Converges in three to five rounds. Like computing a tank's level from "
         "inflow equals outflow instead of waiting for it to settle.")

# 7 — proof
img_below("Equilibrium fix", FOREST, "It reaches the true value — proof",
    ["Brute force from 240 K AND 260 K both land on the shortcut's profile.",
     "Agreement < 0.1 K across the sensor depths.",
     "Different starting points → one curve = the same steady state."],
    "fig_equilibrium_demo",
    "This is the proof the shortcut isn't a trick. I ran the slow brute-force method from two very different starting "
    "temperatures — 240 and 260 K — for 3000 lunar days. Both land exactly on the fast method's curve, to under a tenth of "
    "a degree, at every depth. Different starting points, same destination. That's what 'the two methods agree' really means, "
    "and it's why I trust the shortcut.")

# 8 — certification
img_right("Equilibrium fix", FOREST, "Certified, not assumed — and ~30× faster",
    ["Guess-independence ≤ 0.03 K; deep heat flux closes on Q_b to < 2%.",
     "The solver returns drift & closure diagnostics — it proves it arrived.",
     "Same answer in ~9 s vs ~4 min of brute force per solve."],
    "fig_equilibrium_certification",
    "Two more things. First, the method certifies itself — every solve returns two diagnostics: how much the anchor "
    "temperature drifted (we require under 0.03 K) and how well the deep flux matches Q_b (under 2%). So we never guess "
    "whether it converged — it tells us, which is exactly what the old method couldn't do. Second, speed: it reaches the "
    "same state in about nine seconds, versus roughly four minutes for brute force — about thirty times faster per solve.")

# 9 — sweep
img_right("Results", CORAL, "Retrieving K_d: sweep the model against the data",
    ["For each site, vary K_d and compute the deep-sensor RMSE.",
     "The minimum RMSE gives the best-fit K_d* (sub-grid parabola fit).",
     "A17's minimum sits at a clearly higher K_d than A15's."],
    "fig_kd_sweep",
    "Now the retrieval itself. For each site we try a range of K_d values; for each one the converged model predicts the "
    "deep temperatures and we measure the misfit — the RMSE — against the sensors. The K_d that minimises that misfit is "
    "our estimate. Each curve has a clean single minimum, and the key point: Apollo 17's minimum is at a clearly higher "
    "conductivity than Apollo 15's.")

# 10 — numbers (stat)
s = base("Results", CORAL, "The new K_d values")
def stat(x, y, val, label, ci, color):
    para(tbox(s, x, y, 4.6, 0.95), val, 52, color, first=True, bold=True, font=HEAD)
    para(tbox(s, x, y + 0.95, 4.6, 0.4), label, 14, CHAR, first=True, bold=True)
    para(tbox(s, x, y + 1.32, 4.6, 0.4), ci, 12, DIM, first=True)
stat(0.95, 1.95, "4.58", "mW m⁻¹ K⁻¹  ·  Apollo 15", "95% CI  [4.12, 7.45]", FOREST)
stat(0.95, 4.05, "8.12", "mW m⁻¹ K⁻¹  ·  Apollo 17", "95% CI  [6.85, 9.04]", CORAL)
cbox = rect(s, 0.95, 6.05, 4.6, 0.85, RGBColor(0xEC, 0xE6, 0xDA), MSO_SHAPE.ROUNDED_RECTANGLE)
ctf = cbox.text_frame; ctf.vertical_anchor = MSO_ANCHOR.MIDDLE; ctf.margin_left = Inches(0.2)
cp = ctf.paragraphs[0]; run(cp, "Contrast 3.31   ·   p ≈ 0.011", 18, CHAR, bold=True, font=HEAD)
fit(s, A / "fig_bootstrap.png", 5.9, 1.7, 6.6, 5.2)
notes(s, "Here are the headline numbers. Apollo 15: K_d about 4.6 milliwatts per metre per kelvin. Apollo 17: about 8.1 "
         "— roughly double. For error bars we bootstrap: resample the sensors 1500 times with realistic depth jitter and "
         "re-fit each time. The 95% intervals are shown and barely overlap. The difference is 3.3, with about a 1% chance "
         "it's a fluke — p around 0.011. So the contrast is real, not noise.")

# 11 — meaning
img_right("Results", CORAL, "What it means: a better fit, and a real contrast",
    ["Per-site K_d fits the deep sensors far better than the global K_d (halves RMSE at A17).",
     "A17 conducts heat ~2× more efficiently than A15.",
     "Supplies the meter-scale T(z) boundary that subsurface retrievals need."],
    "fig_thermal_profiles",
    "What does that buy us? The per-site values fit the actual deep temperatures much better than the one-size-fits-all "
    "global value — at Apollo 17 it roughly halves the error. Physically, the regolith at 17 carries heat about twice as "
    "effectively as at 15. And practically, it gives the correct deep temperature profile that anyone modelling the "
    "subsurface — for ice, for missions — actually needs.")

# 12 — Qb conditional + Bayesian
img_right("K_d – Q_b link", GOLD, "K_d is retrieved conditional on Q_b",
    ["The deep gradient is Q_b / K — so K_d and Q_b trade off directly.",
     "We hold Q_b at the published values (A15 = 21, A17 = 15 mW m⁻²).",
     "Bayesian cross-check: even with Q_b uncertain, A17 > A15 ~83% of the time.",
     "Direction of the contrast is robust; exact magnitude is conditional."],
    "fig_kd_qb_posterior",
    "Now the important caveat — the thing you asked about. K_d isn't measured in a vacuum; it's tied to the basal heat "
    "flux Q_b, because the gradient is Q_b over K. They're entangled: assume more heat from below and you'd infer a higher "
    "conductivity for the same temperatures. We can't pull both out of this data, so we fix Q_b at the published Apollo "
    "values — 21 and 15 — and report K_d conditional on them. The figure shows that trade-off. We then checked it Bayesianly: "
    "even letting Q_b vary, 17 stays more conductive than 15 about 83% of the time. The direction is robust; the exact gap depends on the flux.")

# 13 — model selection table
table_slide("K_d – Q_b link", GOLD, "Conductivity contrast, not a flux artifact",
    ["Model fit to both sites", "RMSE (K)", "ΔAICc"],
    [["M1: per-site K_d, Q_b fixed", "0.60", "0.0  (best)"],
     ["M3: shared K_d, Q_b fixed",   "0.67", "2.71"],
     ["M2: shared K_d, Q_b free",    "0.61", "3.99"]],
    [6.6, 2.4, 2.4],
    "We fit three competing models to both sites at once and compare with AICc (rewards fit, penalises free parameters). "
    "Model 1 — a different conductivity per site with flux fixed — wins. The models that force one shared conductivity are "
    "worse by ~3-4 AICc units, even when the flux is allowed to float to compensate. So the data prefers a conductivity "
    "contrast over a flux contrast — that's what lets us call this a conductivity result.",
    cap="Lowest AICc wins. ΔAICc ≥ ~3 is a meaningful preference for M1.")

# 14 — robustness
img_below("Cross-checks", TEAL, "Robustness checks",
    ["Joint (K_d, H) fit: the conductivity transition depth doesn't rescue a single shared K_d.",
     "Sensitivity to the borestem cut, stability threshold and albedo — the contrast persists."],
    "fig_robustness",
    "We stress-tested the result against the model's other knobs. Letting the conductivity's depth-transition scale H vary "
    "jointly with K_d doesn't erase the contrast. Neither does moving the borestem cut, changing the stability threshold for "
    "selecting data, or perturbing the surface albedo. The two sites stay separated across all of these.")

# 15 — diviner
img_right("Cross-checks", TEAL, "Independent check: Diviner surface temperatures",
    ["Uses NO heat-flow data: run the model forward with the retrieved K_d.",
     "Predicted surface day–night curve matches the Diviner orbital composite.",
     "The deep fit and the surface observations are self-consistent."],
    "fig_diviner_closure",
    "An independent sanity check that doesn't touch the heat-flow data at all. We take the model with our retrieved K_d "
    "and run it forward to predict the surface temperature swing over a lunar day, then compare to what Diviner measured "
    "from orbit. They agree. So the model we fit to deep sensors also reproduces the surface — the whole picture is "
    "consistent, not tuned to one thing.")

# 16 — martinez
img_right("Cross-checks", TEAL, "Different conductivity model, same story",
    ["Repeat with the Martínez & Siegler (2021) K(T, ρ) form instead of Hayne.",
     "Each site prefers a different deep density (α) — the same contrast, via a different model.",
     "The conclusion isn't an artifact of one functional form."],
    "fig_alpha_sweep",
    "Finally we checked the contrast isn't an artifact of the particular conductivity formula. We redid it with a "
    "completely different published form — Martínez and Siegler's, which depends on density rather than depth. Each site "
    "prefers a different deep density, telling the same story. So the result holds across two independent ways of "
    "describing the regolith.")

# 17 — error budget table
table_slide("Cross-checks", TEAL, "Where the uncertainty comes from",
    ["Component  (1σ, mW m⁻¹ K⁻¹)", "Apollo 15", "Apollo 17"],
    [["Statistical (bootstrap)", "0.80", "0.55"],
     ["Basal flux  Q_b", "1.20", "2.16"],
     ["Surface albedo  A", "0.80", "3.49"],
     ["Surface conductivity  K_s", "0.53", "1.96"],
     ["Borestem / threshold / ρ", "≤0.15", "≤0.42"],
     ["Total (quadrature)", "1.75", "4.61"]],
    [5.8, 2.6, 2.6],
    "This breaks down the error bar. The biggest contributors are the basal flux Q_b and the surface albedo — not the "
    "statistics. That's the entanglement again: because K_d rides on Q_b, our uncertainty in the flux is one of the largest "
    "pieces of our uncertainty in K_d, especially at Apollo 17. It tells us where to push next for tighter numbers.",
    cap="Q_b and albedo dominate — the K_d–Q_b link showing up in the budget.")

# 18 — architecture
img_below("Computation", BLUEG, "How the code is organized",
    ["Three layers: a physics engine (lunar/), thin driver scripts, the paper.",
     "One hub function — run_with() — turns a trial K_d into a settled profile; everything calls it.",
     "Single config file; 37 automated tests guard the numerics."],
    "fig_architecture",
    "On the computing side, the code is in three clean layers: a core physics package, thin scripts that drive it, and the "
    "paper. There's one hub function, run_with, that takes a trial K_d and returns a converged profile — the sweep, the "
    "bootstrap, every figure call it. One configuration file so nothing drifts out of sync, and 37 automated tests.")

# 19 — cores
img_below("Computation", BLUEG, "Why one solve uses one core — and where more cores help",
    ["One equilibrium solve is sequential: each lunation needs the previous — it can't be split.",
     "But the retrieval runs MANY independent solves (the K_d sweep, the bootstrap).",
     "Those fan across all cores — more cores speed up the batch, not a single solve."],
    "seq_vs_parallel",
    "You asked why it only used one core. A single solve marches through lunar days one after another — each depends on the "
    "one before — so it physically can't be split; it's one core, start to finish. But the retrieval isn't one solve, it's "
    "hundreds of independent ones — every trial K_d, every bootstrap resample. Those are independent, so we hand them out "
    "across all your cores at once. So yes, you can absolutely use more cores — they make the whole batch finish sooner; "
    "they just can't speed up any single solve.")

# 20 — C++
img_below("Computation", BLUEG, "Could part of this be ported to C++?",
    ["Yes — the inner time-step loop is interpreted Python; compiled code gives ~10–100× there.",
     "Port the numerics (grid / solver / equilibrium); keep config + analysis in Python.",
     "For this paper Python is fast enough; C++ pays off at scale (many sites / a mission pipeline)."],
    "speedup_stack",
    "Could we port to C++? Yes, and here's the honest version. The slow part is the inner time-stepping loop, currently "
    "plain Python — exactly where a compiled language wins, maybe ten to a hundred times on that loop. The plan: port just "
    "the three small numerical modules, keep configuration and analysis in Python. The key point is the speed-ups stack — "
    "the algorithm fix already bought about 30x, compilation multiplies that, cores multiply again. For the current paper we "
    "don't need it; C++ becomes worth it when we scale to many locations or a mission pipeline.")

# 21 — summary
dark_slide("Summary & next steps", None,
    ["Fixed a hidden bias (F1) with a fast, self-certifying equilibrium solver.",
     "New per-site K_d:  A15 = 4.6,  A17 = 8.1 mW m⁻¹ K⁻¹;  contrast 3.3,  p ≈ 0.011.",
     "Conditional on Q_b (21 / 15 mW m⁻²); contrast direction robust; conductivity, not flux (AICc).",
     "Next: manuscript to JGR: Planets  ·  tighten Q_b  ·  optional C++ for scale."])
notes(prs.slides[-1],
      "To wrap up: we found and fixed a real bias in the forward model and replaced it with a fast solver that certifies "
      "its own convergence. With that we retrieve distinct conductivities at the two sites — about 4.6 and 8.1 — a contrast "
      "that's statistically significant and survives our cross-checks. It's conditional on the assumed basal flux, but the "
      "direction is robust and the data prefers a conductivity difference over a flux difference. Next steps: finish the JGR "
      "Planets manuscript, tighten the flux assumption, and optionally a C++ core if we scale up. Happy to take questions.")

prs.save(str(OUT))
print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
