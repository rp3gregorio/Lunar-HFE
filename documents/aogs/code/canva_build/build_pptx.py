#!/usr/bin/env python3
"""Build the A0 AOGS poster PPTX from layout.json (native, Canva-editable)."""
import json
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import os; SCRATCH = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'aogs_poster_canva.pptx')
mm = lambda v: Emu(int(round(v * 36000)))
rgb = lambda h: RGBColor.from_string(h.lstrip("#").upper())
ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

spec = json.load(open(f"{SCRATCH}/layout.json"))
prs = Presentation()
prs.slide_width, prs.slide_height = mm(spec["page"][0]), mm(spec["page"][1])
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

def no_line(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False

# layer order = z order: boxes/pills (layer -1) first, content on top
for e in sorted(spec["elements"], key=lambda e: e.get("layer", 0)):
    x, y, w, h = mm(e["x"]), mm(e["y"]), mm(e["w"]), mm(e["h"])
    if e["kind"] == "rect":
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if e.get("rounded") else MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = rgb(e["fill"])
        if e.get("line"):
            shp.shadow.inherit = False
            shp.line.color.rgb = rgb(e["line"])
            shp.line.width = mm(e.get("lw", 0.5))
        else:
            no_line(shp)
        if e.get("rounded"):
            shp.adjustments[0] = 0.06
    elif e["kind"] == "image":
        slide.shapes.add_picture(e["path"], x, y, w, h)
    elif e["kind"] == "table":
        gf = slide.shapes.add_table(len(e["rows"]), len(e["rows"][0]), x, y, w, h)
        tbl = gf.table
        tbl.first_row = False; tbl.horz_banding = False
        for c, cwid in enumerate(e["col_w"]):
            tbl.columns[c].width = mm(cwid)
        for r, row in enumerate(e["rows"]):
            tbl.rows[r].height = mm(e["h"] / len(e["rows"]))
            for c, cell in enumerate(row):
                tc = tbl.cell(r, c)
                tc.fill.solid()
                tc.fill.fore_color.rgb = rgb("#E9EFF2" if r == 0 else "#FFFFFF")
                tc.margin_left = tc.margin_right = mm(1.5)
                tc.margin_top = tc.margin_bottom = mm(0.8)
                tc.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tc.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
                run = p.add_run(); run.text = cell
                f = run.font
                f.name = "Arial"; f.size = Pt(19)
                f.bold = (r == 0 or c == 0); f.color.rgb = rgb("#2A2520")
    else:  # text
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        if e.get("runs"):
            # single-paragraph text with per-run underline (AOGS presenter rule)
            p = tf.paragraphs[0]
            p.alignment = ALIGN[e["align"]]; p.line_spacing = 1.08
            for seg, underline in e["runs"]:
                run = p.add_run(); run.text = seg
                f = run.font
                f.name = "Arial"; f.size = Pt(e["size"])
                f.bold = bool(e.get("bold")); f.italic = bool(e.get("italic"))
                f.underline = bool(underline); f.color.rgb = rgb(e["color"])
        else:
            paras = e["text"].split("\n")
            for i, para in enumerate(paras):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = ALIGN[e["align"]]
                p.line_spacing = 1.08
                if e.get("bullets"):
                    p.space_after = Pt(6)
                run = p.add_run()
                run.text = ("•  " + para) if e.get("bullets") else para
                f = run.font
                f.name = "Arial"; f.size = Pt(e["size"])
                f.bold = bool(e.get("bold")); f.italic = bool(e.get("italic"))
                f.color.rgb = rgb(e["color"])

prs.save(OUT)
print("wrote", OUT)
