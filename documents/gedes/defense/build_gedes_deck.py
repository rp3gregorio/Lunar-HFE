#!/usr/bin/env python3
"""Build the GEDES defense deck ON the official Science Tokyo template.

WHY PYTHON AND NOT pptxgenjs: pptxgenjs cannot open an existing .pptx, so it
cannot inherit a template's masters, layouts, theme fonts or colours. python-pptx
can. This script therefore opens the official template, drops its sample slides,
and rebuilds using ITS OWN layouts -- so the corporate design, the Yu Gothic
theme fonts and the slide-number placeholders all carry through untouched.

THE REQUIREMENT THIS ENFORCES (doctoral-course applicants):
    Part 1  Master's thesis      12 min  + 5 min Q&A
    Part 2  Doctoral plan         6 min  + 7 min Q&A

Every slide carries a `secs` budget. The script sums them per part, compares
against the allocation, prints a timing table, and REFUSES TO WRITE if a part
is over. Backup slides are excluded from the budget -- they exist to be jumped
to during the two Q&A blocks, which is what the 5 + 7 minutes are for.

Edit SLIDES below; the timing follows automatically.

Usage
    python documents/gedes/defense/build_gedes_deck.py
    python documents/gedes/defense/build_gedes_deck.py --template /path/to.pptx
    python documents/gedes/defense/build_gedes_deck.py --allow-over   # report only
"""
from __future__ import annotations

import argparse
import copy
import pathlib
import sys

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = pathlib.Path(__file__).resolve().parent
IMG = HERE / "img"
DEFAULT_TEMPLATE = pathlib.Path(
    "/Users/rp3gregorio/Downloads/パワーポイント_PowerPoint templates_16：9.pptx")
OUT = HERE / "24M58378Gregorio.pptx"

# ---- the template's own layouts, by index (verified against the file) -------
L_TITLE   = 0    # A-タイトルスライド   title / name / affiliation / event
L_DIVIDER = 1    # A-中表紙スライド1    section divider
L_TOC     = 12   # 目次ページ
L_HEAD    = 13   # コンテンツ1-1        title + a heading block, figure below
L_BODY    = 14   # コンテンツ1-2        title + full-height body
L_BLANK   = 21   # 白紙                 slide number only -- best for figures

# allocation from the GEDES instructions, in seconds
BUDGET = {"P1": 12 * 60, "P2": 6 * 60}
QA = {"P1": 5 * 60, "P2": 7 * 60}
# aim a little under, so pauses and an animation loop fit
HEADROOM = 0.96

SITE = dict(
    title="How well does the Moon hold its heat?",
    subtitle="Measuring it at the only two places we ever dug",
    name="Ramon III P. Gregorio",
    affil="Student No. 24M58378 ・ Institute of Science Tokyo ・ Kasai Laboratory\n"
          "Supervisor: Prof. Yasuko Kasai ・ Co-supervisor: Dr. Richard Larsson",
    event="GEDES Master's Thesis Defense ・ July 2026",
)


def S(part, kind, secs, title="", body=None, fig=None, notes="", **kw):
    return dict(part=part, kind=kind, secs=secs, title=title, body=body,
                fig=fig, notes=notes, **kw)


# ---------------------------------------------------------------- the deck
SLIDES = [
    # ============================== PART 1 — master's thesis, 12 min ======
    S("P1", "title", 20, notes="Open with the question, not the filing label."),

    S("P1", "fig", 40, "The surface is violent. A metre down, nothing moves.",
      body="Ice survives only where the ground is cold and steady — and that is "
           "decided below the surface.",
      fig="anim_hook.gif",
      notes="Let the loop run once before speaking. Surface swings 100–390 K "
            "every month; one metre down it does not move."),

    S("P1", "fig", 45, "Every model uses one number for the whole Moon",
      body="It was fitted from orbit, and a satellite only feels the top few centimetres.",
      fig="lay_gap.png",
      notes="Land the two numerals: 3.4 applied everywhere, and ZERO subsurface "
            "measurements that have ever checked it."),

    S("P1", "fig", 55, "Only two holes have ever been drilled and instrumented",
      body="Apollo 15 and 17, 1971–1977, recovered from the original mission tapes.",
      fig="lay_boreholes.png",
      notes="23 sensors survive the 80 cm cut: 7 at Apollo 15, 16 at Apollo 17. "
            "Nothing like it is coming again."),

    S("P1", "fig", 35, "The whole study on one slide",
      body="Three inputs, four steps, three numbers out.",
      fig="gb_pipeline.png",
      notes="Walk the columns with your hand. Only ONE thing in the input column "
            "is unknown. Do not read the equations aloud."),

    S("P1", "fig", 50, "Turning six years of wobble into one honest number",
      body="An automatic rule picks the flattest trailing stretch. Nothing is chosen by hand.",
      fig="lay_window.png",
      notes="Drift limit 0.08 K per year; leftover drift is carried as an error, "
            "not discarded."),

    S("P1", "fig", 40, "A simple picture of the ground, and the equations behind it",
      body="Sunlight in, heat radiated out, a steady trickle from below — and one unknown.",
      fig="lay_model.png",
      notes="Say plainly that this part is standard. That buys credibility for "
            "the next three slides, where something is new."),

    S("P1", "fig", 35, "The model runs blind — the data enters only at the fit",
      body="The forward model never sees a thermometer.",
      fig="gb_dataenters_slide.png",
      notes="This answers 'did you just fit your own data?'. Sets up the "
            "profile slide later."),

    S("P1", "fig", 35, "Why this calculation was impossible",
      body="Three nested loops: ~3000 lunations to settle, ~27 hours per experiment.",
      fig="gb_costnesting.png",
      notes="This is the wall. Not slow — not possible."),

    S("P1", "fig", 70, "The calculation used to take a day. Now it takes a minute.",
      body="Once the ground repeats its monthly cycle, the deep part is rebuilt "
           "from one anchor instead of simulated.",
      fig="lay_solver.png",
      notes="STAR SLIDE. Kettle analogy if eyes glaze. ALWAYS pair the speed "
            "claim with the accuracy claim: same answer to better than 0.01 mW."),

    S("P1", "fig", 65, "Both sites hold heat differently than the textbook value",
      body="Apollo 17 lets heat through about 1.5× more easily, and both exceed "
           "the single global value.",
      fig="lay_results.png",
      notes="STAR SLIDE. 4.60 and 7.08 against a global 3.4. Then the evidence "
            "nobody asks for: the fit got BETTER, 0.89 → 0.40 K at Apollo 17."),

    S("P1", "fig", 50, "The model against the actual thermometers",
      body="Not a number from an optimiser — the profile that passes through "
           "measurements the model never saw.",
      fig="th_profiles.png",
      notes="The pay-off of the blind-model slide. Shallow sensors are drawn "
            "open: excluded there, excluded here."),

    S("P1", "fig", 35, "Re-running the whole analysis 1500 times",
      body="Leaving sensors out at random and jittering their depths gives the "
           "full spread the data supports.",
      fig="lay_bootstrap.png",
      notes="The spreads barely overlap — that is why the ORDERING is solid. "
            "Do NOT say p-value: it is a bootstrap tail proportion, 0.031."),

    S("P1", "fig", 55, "What I can claim, and what I cannot",
      body="A buried thermometer measures steepness, and steepness is "
           "heat-from-below divided by conductivity.",
      fig="lay_seesaw.png",
      notes="STAR SLIDE. Calmly, no apology. Ordering survives; exact size does "
            "not. The 95% range on the difference still touches zero."),

    S("P1", "bullets", 55, "Conclusions",
      body=["The two boreholes are genuinely different — 4.60 and 7.08, about "
            "1.5× apart, and both above the global 3.4",
            "A day-long calculation now takes a minute — ≈2500× faster, which is "
            "what made the uncertainty analysis possible at all",
            "The ordering is robust; the magnitude is not — >99% of tested cases, "
            "but the difference interval still touches zero",
            "This is the ground truth subsurface missions need"],
      notes="Land these, then stop. No thank-you slide."),

    # ============================== PART 2 — doctoral plan, 6 min =========
    S("P2", "divider", 10, "Doctoral Research Plan",
      notes="Change of gear: that is what I have done, here is where it goes."),

    S("P2", "fig", 45, "Five phases, and where the work actually stands",
      body="Two phases are delivered and presented; three are scheduled against them.",
      fig="study_phases_flow.png",
      notes="Point at the dashed rule: left of it exists, right of it is the "
            "proposal. Then stop."),

    S("P2", "fig", 40, "Phase 1 — the thesis, delivered",
      body="Both boreholes reproduced, every claim stress-tested, limits stated honestly.",
      fig="phase1_slide.png",
      notes="Already defended in Part 1 — keep this brief."),

    S("P2", "fig", 60, "Phase 2 — real terrain, at the two sites we can check",
      body="Terrain shifts Apollo 15 down by 2.72 and Apollo 17 up by 2.61 — "
           "opposite directions.",
      fig="phase2_slide.png",
      notes="STRONGEST slide of Part 2. Opposite signs mean no global correction "
            "factor can absorb terrain. That is the argument for Phase 3. "
            "Mention the IR self-heating negative result too."),

    S("P2", "fig", 50, "Phase 3 — from two neighbourhoods to the whole Moon",
      body="Same physics, same solver. What remains is scale.",
      fig="phase3_slide.png",
      notes="Nothing here is a research risk — it is compute and validation. "
            "One column ≈ 1 s is what makes it arithmetic rather than ambition."),

    S("P2", "fig", 45, "Phase 4 — coupling to TSUKIMI: from temperature to ice",
      body="The temperature field becomes terahertz brightness, and finally an "
           "ice-survivability map.",
      fig="phase4_slide.png",
      notes="The 'why me, why here' slide. Name the NICT link explicitly."),

    S("P2", "fig", 35, "Phase 5 — the next-generation regolith model",
      body="Three physics upgrades, each weighed by difficulty and impact before "
           "any is attempted.",
      fig="phase5_slide.png",
      notes="The point is that they were weighed. Do not promise the "
            "vapor-diffusion term."),

    S("P2", "bullets", 45, "Why this is achievable",
      body=["The fast solver exists and is verified — a Moon-wide map is millions "
            "of independent columns, and each now costs a second",
            "The terrain and property studies are finished and presented at AOGS",
            "The home for the work is in place — Institute of Science Tokyo, "
            "Kasai Laboratory, with the NICT terahertz mission link"],
      notes="Then stop for questions. Not a wish list — the continuation of "
            "something already running."),

    # ============================== BACKUP — for the 5 + 7 min Q&A ========
    S("BK", "divider", 0, "Backup",
      notes="Not presented. Jump here by slide number during Q&A."),
]

BACKUP_FIGS = [
    ("Where the uncertainty comes from", "robustness.png",
     "Basal flux dominates at Apollo 15; albedo at Apollo 17.", "why is the error bar so wide?"),
    ("The conductivity–heat-flow trade-off", "qbdeg.png",
     "The two sites respond to a flux revision in opposite directions.",
     "could this be a heat-flow difference instead?"),
    ("Is a per-site value statistically justified?", "aicc.png",
     "Decisive at Apollo 17; at Apollo 15 the global value is not formally rejected.",
     "is the difference significant? Answer honestly."),
    ("Independent check against orbital data", "diviner.png",
     "Surface temperatures were never fitted — a genuine out-of-sample test.",
     "how do you know the model is right?"),
    ("Does it hold when you hold data back?", "holdout.png",
     "TG/TR cross-prediction and leave-one-deepest-out both reproduce the value.",
     "are you over-fitting seven sensors?"),
    ("The Bayesian cross-check", "posterior.png",
     "Floating the basal flux gives the same ordering, at 99.2%.",
     "what if the published heat flow is wrong?"),
    ("The bootstrap, as published", "bootstrap.png",
     "The journal version, with both panels and the tail proportion.",
     "show me the real distribution."),
    ("Finding the answer: the RMSE bowl", "lay_bowl.png",
     "Each site's misfit against candidate conductivity.",
     "how exactly did you pick the value?"),
    ("What the old method had to do", "anim_race.gif",
     "Thousands of cycles against four.", "explain the speed-up."),
    ("The anchor method, step by step", "anchor_method_steps.gif",
     "Step A settles the skin; Step B rebuilds the deep from the anchor.",
     "how does the anchor method actually work?"),
    ("The flux-anchored loop", "gb_fluxanchored.png",
     "Condition, Step A, Step B, and the convergence test.", "show the algorithm."),
    ("From the heat equation to the closure", "gb_pde2odeflow.png",
     "Five moves collapse the PDE to a one-variable ODE.", "where does the closure come from?"),
    ("How one sensor becomes one temperature", "gb_windowflow.png",
     "The stability-window rule as a decision tree.", "how were the windows chosen?"),
    ("Terrain shadowing at both sites", "shadowing.gif",
     "Apollo 15 loses 1.16% of its sunlight to its own horizon; Apollo 17, 0.18%.",
     "how was the DEM used?"),
    ("The AOGS parameter study", "aogs_cpvc.png",
     "The full density sweep behind the doctoral plan.", "show the real analysis."),
    ("The three bugs, and what each cost", "gb_threebugs.png",
     "A wrong basal flux, an under-resolved inner loop, a skipped wrap step.",
     "what went wrong along the way?"),
]


# ------------------------------------------------------------------ build
def strip_slides(prs):
    """Remove the template's sample slides, keeping masters and layouts."""
    ids = prs.slides._sldIdLst
    for sld in list(ids):
        prs.part.drop_rel(sld.rId)
        ids.remove(sld)


def set_ph(slide, idx, text):
    """Fill a layout placeholder, preserving its template formatting."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text_frame.text = text
            return ph
    return None


def drop_ph(slide, idx):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == idx:
            ph._element.getparent().remove(ph._element)
            return


def add_fig(slide, name, top=1.30, bottom=7.05, left=0.30, right=13.03):
    """Place a figure inside the content band, preserving aspect."""
    from PIL import Image
    p = IMG / name
    if not p.exists():
        print(f"    !! missing figure {name}")
        return
    with Image.open(p) as im:
        ar = im.width / im.height
    bw, bh = right - left, bottom - top
    w, h = bw, bw / ar
    if h > bh:
        h, w = bh, bh * ar
    slide.shapes.add_picture(str(p), Inches(left + (bw - w) / 2),
                             Inches(top + (bh - h) / 2), Inches(w), Inches(h))


def add_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def build(template: pathlib.Path, allow_over: bool) -> int:
    prs = Presentation(str(template))
    strip_slides(prs)
    lay = prs.slide_layouts
    rows = []

    def new(layout_idx):
        return prs.slides.add_slide(lay[layout_idx])

    for spec in SLIDES:
        k = spec["kind"]
        if k == "title":
            s = new(L_TITLE)
            set_ph(s, 0, f"{SITE['title']}\n{SITE['subtitle']}")
            set_ph(s, 10, SITE["name"])
            set_ph(s, 11, SITE["affil"])
            set_ph(s, 12, SITE["event"])
        elif k == "divider":
            s = new(L_DIVIDER)
            set_ph(s, 0, spec["title"])
        elif k == "fig":
            s = new(L_HEAD)
            set_ph(s, 0, spec["title"])
            if spec.get("body"):
                ph = set_ph(s, 10, spec["body"])
                if ph is not None:
                    ph.height = Inches(0.62)
                    for p_ in ph.text_frame.paragraphs:
                        for r in p_.runs:
                            r.font.size = Pt(14)
                            r.font.bold = False
            else:
                drop_ph(s, 10)
            add_fig(s, spec["fig"], top=1.78 if spec.get("body") else 1.20)
        elif k == "bullets":
            s = new(L_BODY)
            set_ph(s, 0, spec["title"])
            ph = set_ph(s, 10, "")
            tf = ph.text_frame
            tf.clear()
            # the placeholder is 12.76 in wide, so these lines wrap to two;
            # say so explicitly rather than inheriting it from the layout
            tf.word_wrap = True
            for i, line in enumerate(spec["body"]):
                p_ = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p_.text = line
                p_.level = 0
                for r in p_.runs:
                    r.font.size = Pt(20)
        else:
            raise ValueError(f"unknown kind {k!r}")
        add_notes(s, spec["notes"])
        rows.append((len(prs.slides), spec["part"], spec["secs"], spec["title"] or "Title"))

    for title, fig, take, use in BACKUP_FIGS:
        s = new(L_HEAD)
        set_ph(s, 0, title)
        ph = set_ph(s, 10, take)
        if ph is not None:
            ph.height = Inches(0.62)
            for p_ in ph.text_frame.paragraphs:
                for r in p_.runs:
                    r.font.size = Pt(14)
                    r.font.bold = False
        add_fig(s, fig, top=1.78)
        add_notes(s, f"BACKUP — use if asked: {use}")
        rows.append((len(prs.slides), "BK", 0, title))

    # ------------------------------------------------- timing enforcement
    print(f"\ntemplate : {template.name}")
    print(f"{'#':>3}  {'part':4}  {'sec':>4}  title")
    for n, part, secs, title in rows:
        print(f"{n:3d}  {part:4}  {secs:4d}  {title[:64]}")

    over = False
    print()
    for part in ("P1", "P2"):
        spoken = sum(r[2] for r in rows if r[1] == part)
        cap, target = BUDGET[part], BUDGET[part] * HEADROOM
        n = sum(1 for r in rows if r[1] == part)
        flag = "OVER" if spoken > target else "ok"
        if spoken > target:
            over = True
        print(f"  {part}: {n:2d} slides  {spoken/60:5.2f} min spoken  "
              f"(allocation {cap//60} min, target {target/60:.1f})  "
              f"headroom {(cap-spoken)/60:+.2f} min   [{flag}]")
        print(f"        + {QA[part]//60} min Q&A — covered by the "
              f"{sum(1 for r in rows if r[1]=='BK')} backup slides")

    if over and not allow_over:
        print("\nREFUSING TO WRITE: a part exceeds its target. Trim `secs` in "
              "SLIDES, or pass --allow-over to write anyway.")
        return 1

    prs.save(str(OUT))
    print(f"\nwrote {OUT}  ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", type=pathlib.Path, default=DEFAULT_TEMPLATE)
    ap.add_argument("--allow-over", action="store_true")
    a = ap.parse_args()
    if not a.template.exists():
        sys.exit(f"template not found: {a.template}")
    sys.exit(build(a.template, a.allow_over))
