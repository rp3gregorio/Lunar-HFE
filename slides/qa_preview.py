#!/usr/bin/env python3
"""Layout QA: draw a faithful wireframe of every slide from the real .pptx
shape geometry (exact positions), so overlaps / overflow / misalignment are
visible without LibreOffice. Also flags any shape outside the slide bounds.
"""
import pathlib
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

REPO = pathlib.Path(__file__).resolve().parents[1]
PPTX = REPO / "slides" / "progress_2026-06-17.pptx"
SW, SH = 13.333, 7.5
prs = Presentation(str(PPTX))


def emu_in(v):
    return v / 914400.0

fig, axes = plt.subplots(7, 3, figsize=(16, 22))
flags = []
for idx, slide in enumerate(prs.slides):
    ax = axes[idx // 3][idx % 3]
    ax.set_xlim(0, SW); ax.set_ylim(SH, 0); ax.set_aspect("equal")
    ax.set_title(f"slide {idx+1}", fontsize=9)
    ax.set_xticks([]); ax.set_yticks([])
    # detect background (near full-slide rect) -> use as facecolor, don't draw
    bg = "white"
    for sh in slide.shapes:
        try:
            l, t, w, h = emu_in(sh.left), emu_in(sh.top), emu_in(sh.width), emu_in(sh.height)
        except Exception:
            continue
        if w > 12.5 and h > 6.8:
            bg = "#222"  # dark-ish; just mark it
    ax.set_facecolor("#f0efea" if bg == "white" else "#2a3a30")
    for sh in slide.shapes:
        try:
            l, t, w, h = emu_in(sh.left), emu_in(sh.top), emu_in(sh.width), emu_in(sh.height)
        except Exception:
            continue
        if w > 12.5 and h > 6.8:
            continue  # background, skip
        is_pic = sh.shape_type == MSO_SHAPE_TYPE.PICTURE
        # bounds check
        if l < -0.05 or t < -0.05 or l + w > SW + 0.05 or t + h > SH + 0.05:
            flags.append(f"slide {idx+1}: shape out of bounds ({l:.2f},{t:.2f},{w:.2f}x{h:.2f})")
        ax.add_patch(Rectangle((l, t), w, h, fill=is_pic, facecolor=("#c9d4cd" if is_pic else "none"),
                     edgecolor=("#888" if is_pic else "#3D6E4A"), lw=0.8, alpha=0.9))
        txt = ""
        if sh.has_text_frame:
            txt = sh.text_frame.text.replace("\n", " ⏎ ")
        if is_pic:
            txt = "[IMG]"
        if txt:
            ax.text(l + 0.05, t + 0.18, txt[:46], fontsize=5.5, color="#111", va="top", wrap=True)

# hide unused axes
for k in range(len(list(prs.slides)), 21):
    axes[k // 3][k % 3].axis("off")

fig.tight_layout()
out = REPO / "slides" / "qa_preview.png"
fig.savefig(out, dpi=110, bbox_inches="tight")
print("wrote", out)
print(f"slides: {len(list(prs.slides))}")
if flags:
    print("OUT-OF-BOUNDS FLAGS:")
    for f in flags:
        print("  " + f)
else:
    print("no shapes out of bounds")
